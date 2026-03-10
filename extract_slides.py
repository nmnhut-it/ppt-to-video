"""Extract slide images and text content from a PowerPoint file."""

import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from config import SLIDES_DIR, SCRIPTS_DIR, OUTPUT_DIR


def extract_text_content(pptx_path: str) -> list[dict]:
    """Extract title, body text, and speaker notes from each slide."""
    prs = Presentation(pptx_path)
    slides_data = []

    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_lines = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.shape_id == slide.shapes.title.shape_id if slide.shapes.title else False:
                    title = shape.text_frame.text.strip()
                else:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            body_lines.append(text)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides_data.append({
            "slide_number": i,
            "title": title,
            "body": "\n".join(body_lines),
            "notes": notes,
        })

    return slides_data


def export_slide_images_com(pptx_path: str) -> list[Path]:
    """Export slides as PNG using PowerPoint COM automation."""
    import comtypes.client

    pptx_abs = str(Path(pptx_path).resolve())
    out_dir = str(SLIDES_DIR.resolve())
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)

    ppt = comtypes.client.CreateObject("PowerPoint.Application")
    ppt.Visible = 1
    presentation = ppt.Presentations.Open(pptx_abs, WithWindow=False)

    for i, slide in enumerate(presentation.Slides, start=1):
        out_path = str(Path(out_dir) / f"slide_{i:03d}.png")
        slide.Export(out_path, "PNG", 1920, 1080)

    presentation.Close()
    ppt.Quit()
    return sorted(SLIDES_DIR.glob("*.png"))


def export_slide_images(pptx_path: str) -> list[Path]:
    """Export slides as images. Tries COM (PowerPoint), then placeholder."""
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)

    # Try PowerPoint COM automation
    try:
        images = export_slide_images_com(pptx_path)
        if images:
            print(f"Exported {len(images)} slides via PowerPoint COM")
            return images
    except Exception as e:
        print(f"PowerPoint COM failed: {e}")

    # Fallback: placeholder images
    print("Creating placeholder slide images.")
    prs = Presentation(pptx_path)
    images = []
    width = int(prs.slide_width.inches * 96)
    height = int(prs.slide_height.inches * 96)

    for i, slide in enumerate(prs.slides, start=1):
        img = Image.new("RGB", (width, height), color=(255, 255, 255))
        path = SLIDES_DIR / f"slide_{i:03d}.png"
        img.save(path)
        images.append(path)

    return images


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract_slides.py <path-to-pptx>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    print(f"Extracting from: {pptx_path}")

    # Extract text content
    slides_data = extract_text_content(pptx_path)
    SCRIPTS_DIR.mkdir(parents=True, exist_ok=True)
    content_file = SCRIPTS_DIR / "slides_content.json"
    content_file.write_text(json.dumps(slides_data, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Saved slide content to {content_file} ({len(slides_data)} slides)")

    # Export images
    images = export_slide_images(pptx_path)
    print(f"Exported {len(images)} slide images to {SLIDES_DIR}")


if __name__ == "__main__":
    main()
